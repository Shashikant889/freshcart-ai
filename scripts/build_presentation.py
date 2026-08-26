import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Theme colors
    COLOR_BG = RGBColor(0xFA, 0xFA, 0xFA)
    COLOR_NAVY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_EMERALD = RGBColor(0x04, 0x78, 0x57)
    COLOR_LIGHT_GREEN = RGBColor(0x10, 0xB9, 0x81)
    COLOR_SLATE = RGBColor(0x33, 0x41, 0x55)
    COLOR_MUTED = RGBColor(0x64, 0x74, 0x8B)
    COLOR_CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
    COLOR_BORDER = RGBColor(0xCB, 0xD5, 0xE1)
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    def add_header(slide, title_text, category="FRESHCART AI — MAJOR PROJECT PRESENTATION"):
        # Header banner shape
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_WHITE
        top_bar.line.color.rgb = COLOR_BORDER
        top_bar.line.width = Pt(1)
        
        # Category label
        tf_cat = top_bar.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Calibri'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_EMERALD
        
        # Title
        p_title = tf_cat.add_paragraph()
        p_title.text = title_text
        p_title.font.name = 'Calibri'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_NAVY
        
    def add_card(slide, left, top, width, height, title="", bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        
        if title:
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = 'Calibri'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_NAVY
            p.alignment = PP_ALIGN.LEFT
        return card

    def add_bullet_list(slide, left, top, width, height, items, font_size=16, space_after=12):
        tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = 'Calibri'
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_SLATE
            p.space_after = Pt(space_after)
            p.level = 0
            p.bullet = True
        return tx_box

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_NAVY
    bg1.line.fill.background()
    
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(2.2))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "UNIVERSITY OF MUMBAI  |  B.E. MAJOR PROJECT PRESENTATION (2025–2026)"
    p1.font.name = 'Calibri'
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_LIGHT_GREEN
    p1.space_after = Pt(10)
    
    p2 = tf1.add_paragraph()
    p2.text = "AI-Driven Intelligent Grocery Retail System\nUsing Machine Learning"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_after = Pt(8)
    
    p3 = tf1.add_paragraph()
    p3.text = "An Integrated Quick-Commerce Intelligence Platform combining Predictive ML and Operations Research"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    # Team card
    card_team = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.6), Inches(5.4), Inches(3.2))
    card_team.fill.solid()
    card_team.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card_team.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    tf_team = card_team.text_frame
    tf_team.word_wrap = True
    p_th = tf_team.paragraphs[0]
    p_th.text = "PROJECT GROUP MEMBERS"
    p_th.font.name = 'Calibri'
    p_th.font.size = Pt(13)
    p_th.font.bold = True
    p_th.font.color.rgb = COLOR_LIGHT_GREEN
    p_th.space_after = Pt(8)
    
    members = [
        "Shashikant Shukla (Moodle ID: [STUDENT_1_MOODLE_ID])",
        "Om Dubey (Moodle ID: [STUDENT_2_MOODLE_ID])",
        "Shreyash Wadalkar (Moodle ID: [STUDENT_3_MOODLE_ID])",
        "[STUDENT_4_NAME — DO NOT GUESS] (Moodle ID: [STUDENT_4_MOODLE_ID])"
    ]
    for m in members:
        pm = tf_team.add_paragraph()
        pm.text = f"• {m}"
        pm.font.name = 'Calibri'
        pm.font.size = Pt(12)
        pm.font.color.rgb = COLOR_WHITE
        pm.space_after = Pt(4)
        
    # Institutional card
    card_inst = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(3.6), Inches(5.4), Inches(3.2))
    card_inst.fill.solid()
    card_inst.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card_inst.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    tf_inst = card_inst.text_frame
    tf_inst.word_wrap = True
    p_ih = tf_inst.paragraphs[0]
    p_ih.text = "GUIDANCE & INSTITUTION"
    p_ih.font.name = 'Calibri'
    p_ih.font.size = Pt(13)
    p_ih.font.bold = True
    p_ih.font.color.rgb = COLOR_LIGHT_GREEN
    p_ih.space_after = Pt(8)
    
    inst_info = [
        "Under the Guidance of: [PROJECT_GUIDE_NAME_AND_TITLE]",
        "Department: Computer Science & Engineering (AIML)",
        "Institute: A. P. Shah Institute of Technology (APSIT), Thane",
        "Affiliation: University of Mumbai",
        "Academic Year: 2025–2026"
    ]
    for info in inst_info:
        pi = tf_inst.add_paragraph()
        pi.text = f"• {info}"
        pi.font.name = 'Calibri'
        pi.font.size = Pt(12)
        pi.font.color.rgb = COLOR_WHITE
        pi.space_after = Pt(4)

    # =========================================================================
    # SLIDE 2: PROJECT OVERVIEW
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Project Overview: Executive Summary")
    add_card(s2, 0.8, 1.5, 11.733, 5.4, "Executive Highlights & System Scope")
    add_bullet_list(s2, 1.2, 2.1, 11.0, 4.5, [
        "Unified Intelligent Platform: Couples front-end retail personalization with back-end dark-store picking and last-mile fleet logistics in a single synchronized platform.",
        "Personalized Recommendation: Hybrid Collaborative Filtering + TF-IDF Content Matching (alpha=0.60) yielding F1@10 = 0.5027 and NDCG@10 = 0.9790.",
        "Predictive Demand & Dynamic Pricing: 30-day recursive SARIMAX forecasting (RMSE = 5.83, MAPE = 2.50%) paired with econometric Log-Log OLS price elasticity under [±25%] guardrails.",
        "Real-Time Fraud Risk Scoring: Cost-sensitive Random Forest transaction anomaly detector operating under severe class imbalance (ROC-AUC = 0.6087, zero leakage).",
        "Operations Research Optimization: Continuous review (r, Q) inventory policy (-87.64% cost), dark-store 2D TSP picker routing (-37.48% walk), and CVRP delivery routing (-61.62% distance).",
        "Resilient Microservice Gateway: Node.js Express <-> Python FastAPI architecture featuring a 1500ms circuit breaker with automatic in-process fallback (<25ms p95 latency)."
    ], font_size=15, space_after=10)

    # =========================================================================
    # SLIDE 3: MOTIVATION
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Motivation: Operational Inefficiencies in Grocery Retail")
    
    m_cards = [
        ("Customer Interaction", ["• Catalog browsing cold-start", "• Irrelevant generic search", "• Friction during high-demand checkout"], 0.8, 1.5, 3.6),
        ("Retail Store Operations", ["• 15–25% perishable spoilage", "• Inaccurate static forecasting", "• Unconstrained pricing churn"], 4.8, 1.5, 3.6),
        ("Inventory & Procurement", ["• High holding & setup costs", "• Stockouts on top SKUs", "• Lack of stochastic safety buffers"], 8.8, 1.5, 3.6),
        ("Dark-Store Warehouse", ["• Sequential aisle traversal", "• Excessive picker walking waste", "• Missed 10-min packing SLAs"], 0.8, 4.4, 5.6),
        ("Last-Mile Fleet Logistics", ["• Uncoordinated radial deliveries", "• Low vehicle utilization (<40%)", "• High fuel expenditure & fleet mileage"], 6.8, 4.4, 5.7)
    ]
    for title, bullets, l, t, w in m_cards:
        h = 2.6 if t == 1.5 else 2.5
        add_card(s3, l, t, w, h, title)
        add_bullet_list(s3, l + 0.2, t + 0.6, w - 0.4, h - 0.7, bullets, font_size=13, space_after=4)

    # =========================================================================
    # SLIDE 4: PROBLEM STATEMENT
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Problem Statement")
    card_ps = add_card(s4, 0.8, 1.5, 11.733, 2.2, "Formal Engineering Problem Statement", bg_color=RGBColor(0xFE, 0xF3, 0xC7), border_color=RGBColor(0xF5, 0x9E, 0x0B))
    
    t_ps = s4.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.3))
    tf_ps = t_ps.text_frame
    tf_ps.word_wrap = True
    p_ps = tf_ps.paragraphs[0]
    p_ps.text = '"To design, develop, and benchmark an integrated, resilient, and leak-free AI-driven grocery retail system that combines personalized recommendation, time-series demand forecasting, econometric dynamic pricing, and real-time transaction fraud detection with mathematical inventory, warehouse, and delivery optimization under a fault-tolerant microservice architecture."'
    p_ps.font.name = 'Calibri'
    p_ps.font.size = Pt(17)
    p_ps.font.bold = True
    p_ps.font.italic = True
    p_ps.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
    
    add_card(s4, 0.8, 4.0, 11.733, 2.9, "Key Engineering Challenges Addressed")
    add_bullet_list(s4, 1.2, 4.6, 11.0, 2.1, [
        "Eliminating Siloed Architectures: Bridging the gap between predictive customer models and physical logistics.",
        "Leak-Free Machine Learning: Enforcing strict chronological holdout evaluation and recursive multi-step forecasting.",
        "Constrained Safety Boundaries: Preventing algorithmic instability in dynamic pricing via bounded [±25%] guardrails.",
        "Ultra-Low Latency & High Availability: Achieving sub-25ms p95 latency with automatic circuit breaker fallback."
    ], font_size=14, space_after=6)

    # =========================================================================
    # SLIDE 5: OBJECTIVES & TRACEABILITY
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Engineering Objectives & Traceability Matrix")
    add_card(s5, 0.8, 1.5, 11.733, 5.4, "Mapping Engineering Objectives to Implemented Code Modules")
    
    objs = [
        "OBJ-1 (Personalization): Hybrid CF + CB Engine (alpha=0.60) in recommendation_service.py -> F1@10 = 0.5027, NDCG@10 = 0.9790.",
        "OBJ-2 (Demand Forecast): Recursive SARIMAX(1,1,1)x(1,0,1)7 in demand_service.py -> RMSE = 5.83 units, MAPE = 2.50%.",
        "OBJ-3 (Dynamic Pricing): Log-Log OLS Price Elasticity (Ed) in pricing_service.py -> +22.21% Simulated Revenue Lift (p < 0.001).",
        "OBJ-4 (Fraud Scoring): Cost-Sensitive Random Forest in fraud_service.py -> ROC-AUC = 0.6087 on 1.04% rare fraud.",
        "OBJ-5 (Inventory Policy): Continuous Review (r, Q) EOQ + SS in inventory_opt.py -> -87.64% Cost, 99.88% Service Level.",
        "OBJ-6 (Picker Routing): 2D Euclidean Nearest-Neighbor + 2-Opt TSP in warehouse_opt.py -> -37.48% Walk (0.09% gap).",
        "OBJ-7 (Fleet Routing): Clarke-Wright Savings + 2-Opt CVRP in delivery_opt.py -> -61.62% Travel Distance, 82.9% Util.",
        "OBJ-8 (Resilient Gateway): Node.js AI Gateway (1.5s Circuit Breaker) in services/ai-client.js -> Sub-25ms p95 Latency."
    ]
    add_bullet_list(s5, 1.1, 2.1, 11.2, 4.6, objs, font_size=13, space_after=6)

    # =========================================================================
    # SLIDE 6: LITERATURE SURVEY
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Literature Survey: Summary of Recent IEEE Research (2023–2026)")
    add_card(s6, 0.8, 1.5, 11.733, 5.4, "Comparative Analysis of Surveyed Peer-Reviewed IEEE Publications")
    
    lit_bullets = [
        "[1] Smachylo & Zhuravchak (IEEE CSIT 2024): Hybrid CF + Sentiment analysis boosts precision; limitation: high NLP inference latency.",
        "[2] Bodduluri et al. (IEEE Access 2024): Systematic review of 120+ systems; proves weighted linear hybrids give top stability across sparse catalogs.",
        "[4] Qureshi et al. (IEEE Access 2024): Deep Learning with weather/calendar exogenous regressors; highlights risk of autoregressive lag leakage.",
        "[6] Poongothai et al. (IEEE ICSES 2024): Time-series ML + EOQ cuts stockouts by >40%; lacks stochastic variance modeling for safety stock.",
        "[7] Kumari & Kumar (IEEE InC4 2024): Survey on dynamic pricing trends; unconstrained RL causes customer churn; bounds are essential.",
        "[9] Raut et al. (IEEE OTCON 2024): Random Forest ensemble outperforms single trees on imbalanced POS transaction fraud streams.",
        "[12] Chavan & Nitnaware (IEEE CSNT 2025): Edge computing enables sub-30ms pricing; lacks in-process fallback during edge node crash.",
        "[13] de Assis et al. (IEEE Access 2024): Combinatorial warehouse picking cuts travel >30%; focuses on large plants, not rapid dark stores."
    ]
    add_bullet_list(s6, 1.1, 2.1, 11.2, 4.6, lit_bullets, font_size=13, space_after=6)

    # =========================================================================
    # SLIDE 7: RESEARCH GAP
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Identified Research Gap: Siloed vs. Integrated Retail Systems")
    
    add_card(s7, 0.8, 1.5, 5.6, 5.4, "State of Existing Research (Siloed Approach)", bg_color=RGBColor(0xFE, 0xF2, 0xF2), border_color=RGBColor(0xEF, 0x44, 0x44))
    add_bullet_list(s7, 1.0, 2.2, 5.2, 4.4, [
        "Component Isolation: Algorithms are evaluated in laboratory silos without feedback loops to inventory or dispatch.",
        "Methodological Data Leakage: Temporal shuffling in time-series and deterministic synthetic fraud labels distort metrics.",
        "Unconstrained Dynamic Pricing: Black-box reinforcement learning models produce erratic price swings.",
        "Lack of Fault Tolerance: Systems fail completely when machine learning microservices experience latency spikes.",
        "High Monolithic Latency: Heavy deep learning architectures exceed real-time sub-25ms web application SLAs."
    ], font_size=13, space_after=8)
    
    add_card(s7, 6.9, 1.5, 5.6, 5.4, "FreshCart AI Contribution (Integrated Architecture)", bg_color=RGBColor(0xEC, 0xFD, 0xF5), border_color=RGBColor(0x10, 0xB9, 0x81))
    add_bullet_list(s7, 7.1, 2.2, 5.2, 4.4, [
        "End-to-End Operational Loop: Customer demand shaping directly informs procurement, picking, and delivery dispatch.",
        "Audited Leak-Free Pipelines: Strict chronological holdout splits and recursive multi-step forecasting protocols.",
        "Econometrically Bounded Pricing: Log-Log OLS price elasticity constrained within [±25%] safety guardrails.",
        "Two-Tier Resilient Gateway: Node.js 1.5s circuit breaker with seamless fallback to in-process heuristics.",
        "Ultra-Fast Heuristics: 2D TSP and CVRP algorithms executing in under 3ms on standard multi-core CPUs."
    ], font_size=13, space_after=8)

    # =========================================================================
    # SLIDE 8: PROPOSED SYSTEM ARCHITECTURE
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Proposed High-Level System Architecture (Figure 5.1)")
    
    tiers = [
        ("Client Tier", ["• Customer Storefront PWA", "• Bilingual NLP Search", "• FreshBot Recipe Assistant", "• Admin Management Portal"], 0.8, 1.5, 2.7),
        ("Application Tier", ["• Node.js Express (Port 3000)", "• REST Routing & JWT RBAC", "• ACID Order Processing", "• AI Gateway & Circuit Breaker"], 3.8, 1.5, 2.7),
        ("AI Microservice Tier", ["• Python FastAPI (Port 8000)", "• Singleton Model Registry", "• 4 ML Engines (Recs, Demand, Price, Fraud)", "• 3 OR Solvers (EOQ, TSP, CVRP)"], 6.8, 1.5, 2.7),
        ("Persistence Tier", ["• Relational SQLite Schema", "• 7 Normalized Tables", "• Transaction ACID Locks", "• Interaction Event Logs"], 9.8, 1.5, 2.7)
    ]
    for title, bullets, l, t, w in tiers:
        add_card(s8, l, t, w, 5.4, title)
        add_bullet_list(s8, l + 0.15, t + 0.7, w - 0.3, 4.4, bullets, font_size=13, space_after=8)

    # =========================================================================
    # SLIDE 9: AI INTEGRATION ARCHITECTURE
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "AI Integration & Resilient Fallback Architecture (Figure 5.21)")
    
    add_card(s9, 0.8, 1.5, 11.733, 2.4, "Primary Gateway Flow (Normal Operation: Python FastAPI Online)")
    add_bullet_list(s9, 1.1, 2.1, 11.0, 1.6, [
        "Asynchronous Non-Blocking Dispatch: Express server delegates analytical requests to services/ai-client.js.",
        "REST API Microservice: Python FastAPI daemon (Port 8000) executes pre-warmed in-memory ML models.",
        "Sub-25ms Execution: High-efficiency NumPy/scikit-learn routines return inferences in under 15ms."
    ], font_size=14, space_after=4)
    
    add_card(s9, 0.8, 4.2, 11.733, 2.7, "Resilience & Fallback Flow (Microservice Latency > 1500ms or Offline)", bg_color=RGBColor(0xFE, 0xF3, 0xC7), border_color=RGBColor(0xF5, 0x9E, 0x0B))
    add_bullet_list(s9, 1.1, 4.8, 11.0, 1.9, [
        "Circuit Breaker Trigger: If Python response exceeds 1500ms timeout or connection fails, circuit trips.",
        "In-Process Fallback Execution: Automatically routes to Node.js local JavaScript engines (ml/*.js).",
        "Zero User Interruption: Checkout, recommendations, and inventory replenishment continue without throwing 500 errors."
    ], font_size=14, space_after=4)

    # =========================================================================
    # SLIDE 10: SYSTEM MODULES
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Functional System Modules Breakdown")
    
    mod_cards = [
        ("Customer Storefront", ["• PWA Mobile UI", "• English/Hindi Voice/Text Search", "• Category Filtering & Cart", "• FreshBot Recipe Bundling"], 0.8, 1.5, 2.7),
        ("Admin Portal", ["• Executive Revenue & Order KPIs", "• 30-Day Forecast Visualizer", "• Dynamic Pricing Sandbox", "• Automated Purchase Orders"], 3.8, 1.5, 2.7),
        ("Machine Learning", ["• Hybrid CF+CB Recommendations", "• Recursive SARIMAX Forecaster", "• Log-Log OLS Price Elasticity", "• Random Forest Fraud Classifier"], 6.8, 1.5, 2.7),
        ("Optimization Solvers", ["• Continuous Review (r, Q) EOQ", "• 2D TSP Warehouse Order Picker", "• Clarke-Wright CVRP Fleet Solver", "• In-Process Fallback Hierarchy"], 9.8, 1.5, 2.7)
    ]
    for title, bullets, l, t, w in mod_cards:
        add_card(s10, l, t, w, 5.4, title)
        add_bullet_list(s10, l + 0.15, t + 0.7, w - 0.3, 4.4, bullets, font_size=13, space_after=8)

    # =========================================================================
    # SLIDE 11: RECOMMENDATION SYSTEM
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Personalized Hybrid Recommendation Engine")
    
    add_card(s11, 0.8, 1.5, 6.0, 5.4, "Mathematical Formulation & Pipeline")
    add_bullet_list(s11, 1.0, 2.2, 5.6, 4.4, [
        "User-User Collaborative Filtering (sim_CF): Computes Cosine similarity over implicit user-item interaction vectors.",
        "Content-Based TF-IDF (sim_CB): Measures semantic similarity over product category, sub-category, and ingredient tags.",
        "Weighted Linear Ensemble: S_Hybrid = 0.60 * sim_CF + 0.40 * sim_CB.",
        "Top-K Ranking: Sorts unpurchased items by S_Hybrid and extracts top-10 personalized candidates in 4.86 ms."
    ], font_size=13, space_after=8)
    
    add_card(s11, 7.1, 1.5, 5.4, 5.4, "Holdout Benchmark Results (Table 7.1)", bg_color=RGBColor(0xEC, 0xFD, 0xF5), border_color=RGBColor(0x10, 0xB9, 0x81))
    add_bullet_list(s11, 7.3, 2.2, 5.0, 4.4, [
        "Precision@10: 0.9760 (Accurate Top-10 targeting)",
        "Recall@10: 0.3412 (Broad catalog coverage)",
        "F1-Score@10: 0.5027 (+16.3% over standalone CF)",
        "NDCG@10: 0.9790 (High ranking quality)",
        "Catalog Density Context: High precision reflects dense interaction logs across 31 focused SKUs without train/test leakage."
    ], font_size=13, space_after=8)

    # =========================================================================
    # SLIDE 12: DEMAND FORECASTING
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Time-Series Demand Forecasting (SARIMAX)")
    
    add_card(s12, 0.8, 1.5, 5.2, 5.4, "Model Architecture & Benchmark")
    add_bullet_list(s12, 1.0, 2.2, 4.8, 4.4, [
        "Model Formulation: SARIMAX(1,1,1)x(1,0,1)7 with weekly seasonality (s=7).",
        "Exogenous Regressors: Promotional discount flags and calendar effects (gamma^T X_t).",
        "Recursive Forecasting: Multi-step 30-day projection without lookahead lag leakage.",
        "Evaluation Metrics (Table 7.2):\n  • RMSE = 5.83 units\n  • MAE = 3.89 units\n  • MAPE = 2.50%\n  • Latency = 4.46 ms"
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_1_demand_forecast.png"):
        s12.shapes.add_picture("docs/academic/figures/fig_7_1_demand_forecast.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 13: DYNAMIC PRICING
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Econometric Dynamic Pricing & Elasticity")
    
    add_card(s13, 0.8, 1.5, 5.2, 5.4, "Econometric Model & Simulation")
    add_bullet_list(s13, 1.0, 2.2, 4.8, 4.4, [
        "Log-Log OLS Formulation: ln Q = beta_0 + beta_1 * ln P (Price Elasticity Ed = beta_1).",
        "Empirical Elasticity: Ed = -0.136 (t = -10.21, p < 0.001, R^2 = 0.892).",
        "Safety Guardrails: Hard clipping within [±25%] of base retail price.",
        "Model-Based Simulation (Table 7.3):\n  • Simulated Revenue Lift: +22.21%\n  • Simulated Profit Lift: +58.87%\n  • Note: Model-based simulation under CED assumptions."
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_2_price_elasticity.png"):
        s13.shapes.add_picture("docs/academic/figures/fig_7_2_price_elasticity.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 14: FRAUD DETECTION
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Transaction Fraud Detection & Risk Scoring")
    
    add_card(s14, 0.8, 1.5, 5.2, 5.4, "Cost-Sensitive Random Forest")
    add_bullet_list(s14, 1.0, 2.2, 4.8, 4.4, [
        "Architecture: 100 Decision Trees with balanced class weights.",
        "Input Features: Normalized transaction velocity, basket anomaly ratio, account age.",
        "Holdout Evaluation (Table 7.4):\n  • ROC-AUC = 0.6087\n  • Recall = 0.3864\n  • Precision = 0.0829\n  • F1-Score = 0.1365",
        "Academic Transparency: Modest AUC reflects realistic generalization under 1.04% rare fraud with zero synthetic leakage."
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_3_fraud_roc.png"):
        s14.shapes.add_picture("docs/academic/figures/fig_7_3_fraud_roc.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 15: INVENTORY OPTIMIZATION
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "Continuous Review (r, Q) Inventory Optimization")
    
    add_card(s15, 0.8, 1.5, 5.2, 5.4, "365-Day Simulation Benchmark")
    add_bullet_list(s15, 1.0, 2.2, 4.8, 4.4, [
        "Wilson EOQ: Q* = sqrt(2DS/H) computes optimal batch replenishment size.",
        "Stochastic Safety Stock: SS = Z_0.95 * sqrt(L * sigma_D^2 + D^2 * sigma_L^2).",
        "Reorder Point: ROP = (D * L) + SS triggers automated PO generation.",
        "Simulation Results (Table 7.5):\n  • Total Annual Cost: -87.64% (₹796k -> ₹98k)\n  • Cycle Service Level: 99.88% (vs 75.62%)\n  • Stockout Reduction: -98.31% (890 -> 15 days)"
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_4_inventory_cost.png"):
        s15.shapes.add_picture("docs/academic/figures/fig_7_4_inventory_cost.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 16: WAREHOUSE OPTIMIZATION
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    add_header(s16, "Dark-Store 2D TSP Picker Path Optimization")
    
    add_card(s16, 0.8, 1.5, 5.2, 5.4, "2D TSP Heuristic Solver")
    add_bullet_list(s16, 1.0, 2.2, 4.8, 4.4, [
        "Spatial Modeling: Maps items to (x, y) dark-store Euclidean aisle coordinates.",
        "Two-Phase Algorithm:\n  1. Nearest-Neighbor greedy initialization in O(n^2)\n  2. Intra-tour 2-Opt local search eliminating edge crossings.",
        "Benchmark on 100 Batches (Table 7.6):\n  • Walk Distance: -37.48% (9,685m -> 6,055m)\n  • Optimality Gap: 0.09% (vs exact brute-force)\n  • Solver Latency: 2.34 ms (Near-optimal)"
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_5_warehouse_distance.png"):
        s16.shapes.add_picture("docs/academic/figures/fig_7_5_warehouse_distance.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 17: DELIVERY OPTIMIZATION
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    add_header(s17, "Last-Mile Delivery Fleet Routing (CVRP)")
    
    add_card(s17, 0.8, 1.5, 5.2, 5.4, "Clarke-Wright Savings + 2-Opt")
    add_bullet_list(s17, 1.0, 2.2, 4.8, 4.4, [
        "Formulation: Capacitated Vehicle Routing Problem (CVRP) with Q_veh = 25 kg payload.",
        "Heuristic Solver: Clarke-Wright Savings merges orders by greatest distance saved (s_ij).",
        "Intra-Route 2-Opt: Eliminates loop crossings on each vehicle tour.",
        "Benchmark on 100 Instances (Table 7.7):\n  • Fleet Distance: -61.62% (14,502km -> 5,566km)\n  • Capacity Utilization: 82.9% (vs 38.4%)\n  • Solver Latency: 2.31 ms"
    ], font_size=13, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_6_delivery_routing.png"):
        s17.shapes.add_picture("docs/academic/figures/fig_7_6_delivery_routing.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 18: DATABASE / ER DESIGN
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    add_header(s18, "Relational Database Design & Schema Entities")
    
    db_cards = [
        ("users", ["• id (PK), name, email (UQ)", "• password_hash (bcrypt)", "• role ('customer', 'admin')"], 0.8, 1.5, 3.6),
        ("products", ["• id (PK), name, category", "• price, base_price, stock", "• aisle, rack, shelf, pos_x, pos_y"], 4.8, 1.5, 3.6),
        ("orders & items", ["• id (PK), user_id (FK), total", "• status, fraud_risk_score", "• order_items: product_id, qty"], 8.8, 1.5, 3.6),
        ("sales_history", ["• id (PK), product_id (FK)", "• date, units_sold, is_promo", "• Ingested for SARIMAX"], 0.8, 4.4, 5.6),
        ("user_interactions", ["• id (PK), user_id, product_id", "• interaction_type (view, cart, buy)", "• Ingested for Hybrid CF Recs"], 6.8, 4.4, 5.7)
    ]
    for title, bullets, l, t, w in db_cards:
        h = 2.6 if t == 1.5 else 2.5
        add_card(s18, l, t, w, h, title)
        add_bullet_list(s18, l + 0.2, t + 0.6, w - 0.4, h - 0.7, bullets, font_size=13, space_after=4)

    # =========================================================================
    # SLIDE 19: SECURITY & DATA PROTECTION
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    add_header(s19, "Security Architecture, RBAC & Hardening")
    add_card(s19, 0.8, 1.5, 11.733, 5.4, "Multi-Layer Defensive Security Invariants")
    add_bullet_list(s19, 1.2, 2.1, 11.0, 4.5, [
        "Stateless JWT Authentication: HMAC-SHA256 tokens carrying user claims with expiration timestamps.",
        "Role-Based Access Control (RBAC): Protected admin and dispatch endpoints restricted via middleware/auth.js.",
        "Bcrypt Password Hashing: Secure password hashing with salt rounds preventing credential exposure.",
        "SQL Injection Immunity: 100% parameterized SQLite prepared statements across all data access routines.",
        "Input Sanitization & DTO Validation: Strict body size limits (2MB) and schema type checking on all API routes.",
        "Centralized Error Masking: Production-mode error handling prevents stack trace leakage to client browsers."
    ], font_size=15, space_after=10)

    # =========================================================================
    # SLIDE 20: EXPERIMENTAL SETUP & METHODOLOGY
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    add_header(s20, "Experimental Environment & Data Provenance")
    
    add_card(s20, 0.8, 1.5, 5.6, 5.4, "Software & Execution Stack (Table 6.2)")
    add_bullet_list(s20, 1.0, 2.2, 5.2, 4.4, [
        "Runtime: Node.js v20.x, Python v3.12 (FastAPI v0.111.0).",
        "Persistence: SQLite3 / sql.js WebAssembly.",
        "Data Science: NumPy v1.26.4, SciPy v1.13.1, Pandas v2.2.2.",
        "Machine Learning: Scikit-Learn v1.5.0, Statsmodels v0.14.2.",
        "Testing Harness: Node.js Multi-Tier Regression Suite."
    ], font_size=13, space_after=8)
    
    add_card(s20, 6.9, 1.5, 5.6, 5.4, "Dataset Taxonomy & Provenance (Table 6.3)")
    add_bullet_list(s20, 7.1, 2.2, 5.2, 4.4, [
        "Product Catalog: 31 Seeded Grocery SKUs across 4 aisles.",
        "Sales History: 11,315 Daily SKU Sales Records (365 Days).",
        "User Interactions: 83,760 Interaction Events (Views/Cart/Buy).",
        "Customer Personas: 100 Calibrated Synthetic Profiles.",
        "Order Transactions: 4,231 Orders (1.04% rare fraud rate)."
    ], font_size=13, space_after=8)

    # =========================================================================
    # SLIDE 21: MASTER ML RESULTS TABLE
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    add_header(s21, "Consolidated Machine Learning Holdout Results")
    add_card(s21, 0.8, 1.5, 11.733, 5.4, "Empirical Evaluation Metrics across the 4 ML Engines")
    
    ml_results = [
        "Hybrid Recommendations: F1@10 = 0.5027 | NDCG@10 = 0.9790 | Precision@10 = 0.9760 | Recall@10 = 0.3412 (4.86 ms).",
        "30-Day Demand Forecasting: RMSE = 5.83 units | MAE = 3.89 units | MAPE = 2.50% (SARIMAX with promo regressors, 4.46 ms).",
        "Dynamic Price Elasticity: Ed = -0.136 (p < 0.001) | Simulated Revenue Lift = +22.21% | Simulated Profit Lift = +58.87% (CED Model).",
        "Transaction Fraud Risk: ROC-AUC = 0.6087 | Recall = 0.3864 | Precision = 0.0829 | F1 = 0.1365 (0% synthetic leakage, 19.77 ms).",
        "Methodological Integrity: All models evaluated on leak-free chronological holdout partitions without train/test contamination."
    ]
    add_bullet_list(s21, 1.1, 2.2, 11.2, 4.4, ml_results, font_size=14, space_after=10)

    # =========================================================================
    # SLIDE 22: MASTER OPTIMIZATION RESULTS TABLE
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    add_header(s22, "Operations Research Optimization Benchmarks")
    add_card(s22, 0.8, 1.5, 11.733, 5.4, "Empirical Performance across the 3 OR Solvers")
    
    or_results = [
        "Continuous Review (r, Q) Inventory: -87.64% Cost Reduction (₹796k -> ₹98k) | 99.88% Service Level | -98.31% Stockouts (365-day sim).",
        "Dark Store 2D TSP Picker: -37.48% Walking Distance (9,685m -> 6,055m) | 0.09% Optimality Gap vs. Exact Brute Force | 2.34 ms Latency.",
        "Last-Mile Delivery CVRP: -61.62% Fleet Distance (14,502km -> 5,566km) | 82.9% Vehicle Capacity Utilization | 2.31 ms Latency.",
        "Computational Feasibility: All combinatorial solvers converge in under 3 milliseconds on standard multi-core hardware."
    ]
    add_bullet_list(s22, 1.1, 2.2, 11.2, 4.4, or_results, font_size=15, space_after=12)

    # =========================================================================
    # SLIDE 23: SYSTEM PERFORMANCE & LATENCY
    # =========================================================================
    s23 = prs.slides.add_slide(blank_layout)
    add_header(s23, "System Performance, API Latency & Code QA")
    
    add_card(s23, 0.8, 1.5, 5.2, 5.4, "Gateway Latency & Verification")
    add_bullet_list(s23, 1.0, 2.2, 4.8, 4.4, [
        "Local Development Latency (p95):\n  • Catalog Listing: 3.67 ms\n  • Top-K Recs: 7.90 ms\n  • Demand Forecast: 8.80 ms\n  • Dynamic Pricing: 9.87 ms\n  • 2D TSP Picker: 4.40 ms\n  • CVRP Dispatch: 10.83 ms\n  • Fraud Scoring: 19.77 ms",
        "Automated Test Coverage:\n  • 113/113 Assertions Passed (100%)\n  • 56/56 Master Code Checks Passed"
    ], font_size=12.5, space_after=6)
    
    if os.path.exists("docs/academic/figures/fig_7_7_latency_benchmark.png"):
        s23.shapes.add_picture("docs/academic/figures/fig_7_7_latency_benchmark.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # =========================================================================
    # SLIDE 24: APPLICATION DEMO VIEWS
    # =========================================================================
    s24 = prs.slides.add_slide(blank_layout)
    add_header(s24, "Application Demonstration: Storefront & Admin Portal")
    
    add_card(s24, 0.8, 1.5, 5.7, 5.4, "Customer Storefront PWA (Port 3000)")
    add_bullet_list(s24, 1.0, 2.0, 5.3, 1.4, [
        "Instant catalog search across 31 SKUs",
        "Personalized Top-K Hybrid Recommendations",
        "Cart assembly & real-time pricing"
    ], font_size=12, space_after=3)
    if os.path.exists("docs/academic/screenshots/SHOT-01-storefront.png"):
        s24.shapes.add_picture("docs/academic/screenshots/SHOT-01-storefront.png", Inches(1.0), Inches(3.6), width=Inches(5.3))
        
    add_card(s24, 6.8, 1.5, 5.7, 5.4, "Admin Analytics & ML Management (Port 3000)")
    add_bullet_list(s24, 7.0, 2.0, 5.3, 1.4, [
        "Executive KPIs & 30-Day Sales Trend",
        "AI Demand Forecasting & Dynamic Pricing",
        "Automated Stock Alerts & VRP Dispatch"
    ], font_size=12, space_after=3)
    if os.path.exists("docs/academic/screenshots/SHOT-07-admin-dashboard.png"):
        s24.shapes.add_picture("docs/academic/screenshots/SHOT-07-admin-dashboard.png", Inches(7.0), Inches(3.6), width=Inches(5.3))

    # =========================================================================
    # SLIDE 25: END-TO-END OPERATIONAL WORKFLOW
    # =========================================================================
    s25 = prs.slides.add_slide(blank_layout)
    add_header(s25, "End-to-End Operational Lifecycle Workflow")
    add_card(s25, 0.8, 1.5, 11.733, 5.4, "Synchronized Feedback Loop across Customer & Operational Tiers")
    
    e2e_steps = [
        "1. Customer Browsing & Interaction: User interactions feed the Hybrid Recommendation Engine to personalize product displays.",
        "2. Cart Assembly & Dynamic Pricing: Bounded Log-Log OLS pricing optimizes revenue while preventing consumer churn.",
        "3. Atomic Checkout & Fraud Scoring: Transactions are risk-scored in <20ms; valid orders decrement stock atomically in SQLite.",
        "4. Demand Forecasting & Procurement: Daily sales update SARIMAX models; Continuous Review (r, Q) triggers automated POs when stock <= ROP.",
        "5. Dark-Store Warehouse Picking: Placed orders generate 2D TSP pick paths, cutting picker walking distance by 37.48%.",
        "6. Last-Mile Vehicle Dispatch: Assembled orders are clustered into capacity-constrained CVRP delivery routes, cutting travel miles by 61.62%."
    ]
    add_bullet_list(s25, 1.1, 2.2, 11.2, 4.4, e2e_steps, font_size=14, space_after=8)

    # =========================================================================
    # SLIDE 26: PROJECT CONTRIBUTIONS
    # =========================================================================
    s26 = prs.slides.add_slide(blank_layout)
    add_header(s26, "Engineering Contributions vs. Standard Algorithms")
    
    add_card(s26, 0.8, 1.5, 5.6, 5.4, "Standard Foundational Algorithms", bg_color=RGBColor(0xF1, 0xF5, 0xF9))
    add_bullet_list(s26, 1.0, 2.2, 5.2, 4.4, [
        "Cosine Similarity & TF-IDF Vectorization",
        "SARIMAX State-Space Time-Series Modeling",
        "Ordinary Least Squares (OLS) Linear Regression",
        "Random Forest Classification Ensemble",
        "Wilson Economic Order Quantity (EOQ)",
        "2-Opt Local Search Heuristic",
        "Clarke-Wright Savings Algorithm"
    ], font_size=14, space_after=8)
    
    add_card(s26, 6.9, 1.5, 5.6, 5.4, "Project Engineering Contributions", bg_color=RGBColor(0xEC, 0xFD, 0xF5), border_color=RGBColor(0x10, 0xB9, 0x81))
    add_bullet_list(s26, 7.1, 2.2, 5.2, 4.4, [
        "Unified Multi-Tier Synergy: Coupling front-end demand with back-end dark-store picking and fleet dispatch.",
        "Resilient AI Gateway: 1.5s circuit breaker with seamless fallback to in-process Node.js heuristics.",
        "Leak-Free Validation: Rigorous chronological holdout benchmarking across all ML pipelines.",
        "Rapid Heuristic Solvers: Sub-3ms execution for combinatorial warehouse and fleet routing."
    ], font_size=14, space_after=8)

    # =========================================================================
    # SLIDE 27: LIMITATIONS & CONSTRAINTS
    # =========================================================================
    s27 = prs.slides.add_slide(blank_layout)
    add_header(s27, "Academic Limitations & Technical Constraints")
    add_card(s27, 0.8, 1.5, 11.733, 5.4, "Explicitly Acknowledged Constraints (Examiner Transparency)")
    add_bullet_list(s27, 1.2, 2.1, 11.0, 4.5, [
        "Synthetic Data Calibration: Evaluated on calibrated synthetic sales streams; real commercial deployments exhibit higher noise and extreme market shocks.",
        "Model-Based Pricing Simulation: +22.21% revenue lift is an empirical simulation under Constant Elasticity of Demand (CED) assumptions.",
        "Imbalanced Fraud Baseline: ROC-AUC of 0.6087 represents an operational screening pipeline rather than an autonomous production-grade prevention engine.",
        "Local Development Benchmarking: Sub-25ms latency measured on a multi-core local host; cloud hops may introduce minor network overhead.",
        "Static Dark-Store Layout: Assumes unobstructed 2D grid aisles without real-time physical human or robot congestion modeling."
    ], font_size=14, space_after=8)

    # =========================================================================
    # SLIDE 28: FUTURE RESEARCH DIRECTIONS
    # =========================================================================
    s28 = prs.slides.add_slide(blank_layout)
    add_header(s28, "Future Research & Development Scope")
    add_card(s28, 0.8, 1.5, 11.733, 5.4, "Roadmap for Next-Generation Platform Evolution")
    add_bullet_list(s28, 1.2, 2.1, 11.0, 4.5, [
        "Multi-Agent Deep Reinforcement Learning for Pricing: Modeling competitive multi-seller game dynamics in real-time markets.",
        "Automated Mobile Robot (AMR) Warehouse Swarms: Orchestrating physical automated robotic pickers in high-density automated dark stores.",
        "Dynamic Real-Time CVRP with Live Traffic: Integrating GPS telemetry and road congestion APIs to dynamically re-route fleets in transit.",
        "Cloud-Native Kubernetes Deployment: Containerizing the two-tier architecture for automated horizontal pod autoscaling under enterprise-scale traffic.",
        "Graph Neural Network (GNN) Recommendations: Capturing complex multi-modal user-item graph relationships across massive retail catalogs."
    ], font_size=14, space_after=8)

    # =========================================================================
    # SLIDE 29: CONCLUSION
    # =========================================================================
    s29 = prs.slides.add_slide(blank_layout)
    add_header(s29, "Conclusion: Summary of Findings")
    add_card(s29, 0.8, 1.5, 11.733, 5.4, "Fulfillment of Project Objectives & Key Takeaways")
    add_bullet_list(s29, 1.2, 2.1, 11.0, 4.5, [
        "Successfully Built & Integrated: Developed a complete full-stack intelligent grocery retail system uniting 4 ML models and 3 OR optimizers.",
        "Rigorous Empirical Verification: Evaluated on leak-free holdout datasets and 100 combinatorial benchmark instances.",
        "High Operational Resilience: Designed an asynchronous AI Gateway with a 1.5s circuit breaker and zero-crash in-process fallback.",
        "Sub-25ms Real-Time Performance: Met web application latency targets across all analytical and transactional endpoints.",
        "Defensible Engineering Artifact: Fully verified with 113/113 passed test assertions, 56/56 master audit checks, and locked IEEE references."
    ], font_size=15, space_after=10)

    # =========================================================================
    # SLIDE 30: REFERENCES & THANK YOU
    # =========================================================================
    s30 = prs.slides.add_slide(blank_layout)
    bg30 = s30.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg30.fill.solid()
    bg30.fill.fore_color.rgb = COLOR_NAVY
    bg30.line.fill.background()
    
    t_box30 = s30.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.5))
    tf30 = t_box30.text_frame
    tf30.word_wrap = True
    p30_1 = tf30.paragraphs[0]
    p30_1.text = "KEY IEEE REFERENCES & FINAL DEFENSE"
    p30_1.font.name = 'Calibri'
    p30_1.font.size = Pt(12)
    p30_1.font.bold = True
    p30_1.font.color.rgb = COLOR_LIGHT_GREEN
    p30_1.space_after = Pt(4)
    
    p30_2 = tf30.add_paragraph()
    p30_2.text = "Selected IEEE Xplore Citations (2023–2026)"
    p30_2.font.name = 'Calibri'
    p30_2.font.size = Pt(22)
    p30_2.font.bold = True
    p30_2.font.color.rgb = COLOR_WHITE
    
    # Reference Box
    card_ref = s30.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.8))
    card_ref.fill.solid()
    card_ref.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card_ref.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    tf_r = card_ref.text_frame
    tf_r.word_wrap = True
    
    selected_refs = [
        "[1] Smachylo & Zhuravchak, 'Enhancing Recommender Systems,' IEEE CSIT, 2024, doi: 10.1109/CSIT65290.2024.10982556.",
        "[4] Qureshi et al., 'Demand Forecasting in SCM for Rossmann Stores,' IEEE Access, 2024, doi: 10.1109/ACCESS.2024.3472499.",
        "[8] Karunakaran et al., 'Integrating AI & ML for Dynamic Pricing,' IEEE ICPECTS, 2024, doi: 10.1109/ICPECTS62210.2024.10780375.",
        "[10] Mienye & Jere, 'Deep Learning for Credit Card Fraud Detection,' IEEE Access, 2024, doi: 10.1109/ACCESS.2024.3426955.",
        "[13] de Assis et al., 'Optimising Warehouse Order Picking,' IEEE Access, 2024, doi: 10.1109/ACCESS.2024.3497592.",
        "[14] Nugroho & Girsang, 'Three-Layer Multi-Objective VRP Solver with 2-opt,' IEEE ICE3IS, 2025."
    ]
    for idx, r in enumerate(selected_refs):
        p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
        p.text = r
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p.space_after = Pt(2)
        
    # Thank You Banner
    card_ty = s30.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.3), Inches(11.333), Inches(1.6))
    card_ty.fill.solid()
    card_ty.fill.fore_color.rgb = COLOR_EMERALD
    card_ty.line.fill.background()
    tf_ty = card_ty.text_frame
    tf_ty.word_wrap = True
    pty1 = tf_ty.paragraphs[0]
    pty1.text = "THANK YOU!"
    pty1.font.name = 'Calibri'
    pty1.font.size = Pt(26)
    pty1.font.bold = True
    pty1.font.color.rgb = COLOR_WHITE
    pty1.alignment = PP_ALIGN.CENTER
    
    pty2 = tf_ty.add_paragraph()
    pty2.text = "Open for Questions & Viva Examination"
    pty2.font.name = 'Calibri'
    pty2.font.size = Pt(14)
    pty2.font.color.rgb = COLOR_WHITE
    pty2.alignment = PP_ALIGN.CENTER

    output_path = "docs/academic/FINAL_PROJECT_PRESENTATION.pptx"
    prs.save(output_path)
    print(f"Successfully generated 30-slide presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
